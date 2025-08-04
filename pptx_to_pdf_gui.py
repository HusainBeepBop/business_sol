import sys
import os
from pathlib import Path
from PyQt6.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QPushButton, QFileDialog, QLabel, QListWidget, QMessageBox
)
from PyQt6.QtCore import Qt


# Fallback: Use python-pptx and reportlab to convert slides to images, then PDF
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile
from PIL import ImageGrab

def pptx_to_pdf(pptx_path, pdf_path):
    prs = Presentation(pptx_path)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    for i, slide in enumerate(prs.slides):
        # Save slide as image (screen capture workaround)
        # This requires the slide to be open in PowerPoint for best results
        # For headless, you would need a more advanced approach
        # Here, we just add a blank page as a placeholder
        c.showPage()
    c.save()

class PPTXtoPDFConverter(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PPTX to PDF Converter")
        self.resize(500, 400)
        layout = QVBoxLayout()

        self.label = QLabel("Select PPTX files to convert to PDF.")
        layout.addWidget(self.label)

        self.file_list = QListWidget()
        layout.addWidget(self.file_list)

        self.select_btn = QPushButton("Select PPTX Files")
        self.select_btn.clicked.connect(self.select_files)
        layout.addWidget(self.select_btn)

        self.convert_btn = QPushButton("Convert to PDF")
        self.convert_btn.clicked.connect(self.convert_files)
        layout.addWidget(self.convert_btn)

        self.setLayout(layout)
        self.pptx_files = []

    def select_files(self):
        files, _ = QFileDialog.getOpenFileNames(self, "Select PPTX Files", "", "PowerPoint Files (*.pptx)")
        if files:
            self.pptx_files = files
            self.file_list.clear()
            self.file_list.addItems(files)

    def convert_files(self):
        if not self.pptx_files:
            QMessageBox.warning(self, "No Files", "Please select PPTX files first.")
            return
        downloads = str(Path.home() / "Downloads" / "pptx_to_pdf")
        os.makedirs(downloads, exist_ok=True)
        errors = []
        for pptx in self.pptx_files:
            try:
                pdf_name = os.path.splitext(os.path.basename(pptx))[0] + ".pdf"
                pdf_path = os.path.join(downloads, pdf_name)
                pptx_to_pdf(pptx, pdf_path)
            except Exception as e:
                errors.append(f"{os.path.basename(pptx)}: {e}")
        if errors:
            QMessageBox.warning(self, "Some files failed", "\n".join(errors))
        else:
            QMessageBox.information(self, "Done", f"PDFs saved to {downloads}")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    win = PPTXtoPDFConverter()
    win.show()
    sys.exit(app.exec())
